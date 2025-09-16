from tkinter import filedialog
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import comtypes.client
import os

# Load CSV
team = "ORE_BEA"
csv_file = filedialog.askopenfilename()
df = pd.read_csv(csv_file)

# Filter to team pitchers
df = df[df['PitcherTeam'] == team]

# Clean or prepare data
df = df[df['TaggedPitchType'].notna()]  # Remove rows without valid pitches

# Optional: normalize PitchCall names
df['PitchCall'] = df['PitchCall'].fillna("Undefined")

# Keep only pitch calls relevant to BB/K/IP
valid_calls = ['BallCalled', 'StrikeCalled', 'StrikeSwinging', 'FoulBall', 'InPlay', 'HitByPitch']
df = df[df['PitchCall'].isin(valid_calls)]

hit_results = ['Single', 'Double', 'Triple', 'HomeRun']
df['isHit'] = df['PlayResult'].isin(hit_results)

# Tag outcomes
df['isK'] = df['KorBB'] == "Strikeout"
df['isBB'] = df['KorBB'] == "Walk"
df['isIP'] = df['PitchCall'] == "InPlay"

# Estimate innings pitched (3 batters per inning from InPlay)
def estimate_ip(group):
    in_play = group['PitchCall'] == 'InPlay'
    hbp = group['PitchCall'] == 'HitByPitch'
    strikeout = group['KorBB'] == 'Strikeout'
    walk = group['KorBB'] == 'Walk'
    
    batters_faced = (in_play | hbp | strikeout | walk).sum()
    return batters_faced / 3.0

df['isHR'] = df['PlayResult'] == 'HomeRun'

def to_baseball_innings(ip):
    full = int(ip)
    frac = ip - full
    if frac < 1/3:
        return full + 0.0
    elif frac < 2/3:
        return full + 0.1
    else:
        return full + 0.2

# Aggregate by pitcher
pitcher_stats = df.groupby('Pitcher').agg({
    'isK': 'sum',
    'isBB': 'sum',
    'isHR': 'sum',
    'isHit': 'sum'
}).reset_index()

# Calculate IP for each pitcher
pitcher_stats['IP'] = df.groupby('Pitcher').apply(estimate_ip).values
pitcher_stats['TotalPitches'] = df.groupby('Pitcher').size().values

# Rename columns
pitcher_stats.columns = ['Pitcher', 'K', 'BB', 'HR', 'H', 'IP', 'TotalPitches']

# FIP and WAR estimation
FIP_constant = 3.75
replacement_fip = 7.44

pitcher_stats['WHIP'] = (pitcher_stats['BB'] + pitcher_stats['H']) / pitcher_stats['IP']
pitcher_stats['FIP'] = (13 * pitcher_stats['HR']+ 3 * pitcher_stats['BB'] - 2 * pitcher_stats['K']) / pitcher_stats['IP'] + FIP_constant
pitcher_stats['Runs_Prevented'] = pitcher_stats['IP'] * (replacement_fip - pitcher_stats['FIP']) / 9
pitcher_stats['WAR'] = pitcher_stats['Runs_Prevented'] / 9

# Sort by WAR
pitcher_stats['IP'] = pitcher_stats['IP'].apply(to_baseball_innings)
pitcher_stats['K'] = pitcher_stats['K'].astype(int)
pitcher_stats['BB'] = pitcher_stats['BB'].astype(int)
pitcher_stats['HR'] = pitcher_stats['HR'].astype(int)
pitcher_stats['FIP'] = pitcher_stats['FIP'].round(2)
pitcher_stats['WHIP'] = pitcher_stats['WHIP'].round(2)
pitcher_stats['WAR'] = pitcher_stats['WAR'].round(2)
pitcher_stats = pitcher_stats.sort_values(by='WAR', ascending=False)

# Format pitcher names to "Last F."
def _format_last_first_initial(name: str) -> str:
    try:
        parts = str(name).split(',')
        last = parts[0].strip()
        first_initial = parts[1].strip()[0] + '.' if len(parts) > 1 and len(parts[1].strip()) > 0 else ''
        if first_initial:
            return f"{last}, {first_initial}".strip()
        return last
    except Exception:
        return str(name)

pitcher_stats['Pitcher'] = pitcher_stats['Pitcher'].astype(str).apply(_format_last_first_initial)

# Display final stats
display_cols = ['Pitcher', 'IP', 'K', 'BB', 'HR', 'FIP', 'WHIP', 'WAR']
#print(pitcher_stats[display_cols])


def export_to_pptx(df_results, template_path, output_path):
    df_results = df_results.sort_values(by='WAR', ascending=False).reset_index(drop=True)
    # Load template and use the first slide
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Compute table dimensions with small margins
    emu_per_inch = 914400
    slide_width_inches = prs.slide_width / emu_per_inch
    slide_height_inches = prs.slide_height / emu_per_inch
    margin_h = 0.25
    margin_v = 0.25
    table_width = slide_width_inches - (2 * margin_h)
    table_height = slide_height_inches - (2 * margin_v)
    left = Inches(margin_h)
    top = Inches(margin_v)

    # Restrict to display columns in this order if present
    cols = [c for c in display_cols if c in df_results.columns]
    df_table = df_results[cols]

    rows, cols_n = df_table.shape[0] + 1, df_table.shape[1]
    table = slide.shapes.add_table(rows, cols_n, left, top, Inches(table_width), Inches(table_height)).table

    # Even column widths
    col_width = table_width / cols_n
    for i in range(cols_n):
        table.columns[i].width = Inches(col_width)

    # Header styling
    for j, col_name in enumerate(df_table.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.name = 'Bahnschrift'
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.size = Pt(9)
        para.font.bold = True

    # Body rows
    row_height = table_height / rows
    for row_idx, (_, row) in enumerate(df_table.iterrows()):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.name = 'Bahnschrift'
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.font.size = Pt(9)
        table.rows[row_idx + 1].height = Inches(row_height)

    # Robust save with fallback filename to avoid file-lock issues
    actual_output_path = output_path
    try:
        prs.save(actual_output_path)
    except PermissionError:
        try:
            if os.path.exists(actual_output_path):
                os.remove(actual_output_path)
            prs.save(actual_output_path)
        except Exception:
            base, ext = os.path.splitext(output_path)
            counter = 1
            while True:
                candidate = f"{base}-{counter}{ext}"
                if not os.path.exists(candidate):
                    prs.save(candidate)
                    actual_output_path = candidate
                    break
                counter += 1

    return actual_output_path


def pptx_to_pdf(inputFileName, outputFileName):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, 32)
    deck.Close()
    powerpoint.Quit()
    print(f"Converted {inputFileName} 	 {outputFileName}")


# Export to PPTX and PDF using pitcher template
pptx_path = export_to_pptx(pitcher_stats, "pitcher-template.pptx", "pitcher-results.pptx")
pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
pptx_to_pdf(pptx_path, pdf_path)