from pptx import Presentation, util, text
from pptx.util import Inches, Pt
from tkinter import filedialog
import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
import re
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from tabulate import tabulate
import os
import comtypes.client
from pptx.dml.color import RGBColor

team = ("ORE_BEA")
csv_file =filedialog.askopenfilename()
df = pd.read_csv(csv_file)

#total_runs = df['RunsScored'].sum()
#runs_per_pa = total_runs / 273516
#numerator = (
#    0.689 * df['KorBB'].eq('BB').sum() +
#    0.720 * (df['PitchCall'] == 'HitByPitch').sum() +
#    0.844 * (df['PlayResult'] == 'Single').sum() +
#    1.261 * (df['PlayResult'] == 'Double').sum() +
#    1.601 * (df['PlayResult'] == 'Triple').sum() +
#    2.072 * (df['PlayResult'] == 'HomeRun').sum()
#)
#new_woba = numerator / 273516
#scale = new_woba / runs_per_pa
#print(scale)
#AVG OPS = .785
#AVG SLG = .418
#PA = 273516

# Define constants
W_OBA_WEIGHTS = {'BB': 0.69, 'HBP': 0.72, '1B': 0.89, '2B': 1.27, '3B': 1.62, 'HR': 2.10}
W_OBA_SCALE = 1.264
LEAGUE_WOBA = 0.276
RUNS_PER_WIN = 8.734
LEAGUE_OBP = 0.367  
LEAGUE_SLG = 0.418

# Map PlayResult to hit outcomes
hit_map = {
    'Single': '1B', 'Double': '2B', 'Triple': '3B', 'HomeRun': 'HR'
}
df['HitType'] = df['PlayResult'].map(hit_map)
df['BB'] = df['KorBB'].apply(lambda x: 1 if x == 'Walk' else 0)
df['HBP'] = df['PitchCall'].apply(lambda x: 1 if x == 'HitByPitch' else 0)

# Filter hitters by selected team
team_hitters_df = df[df['BatterTeam'] == team].copy()

# Get total plate appearances
team_hitters_df['PA_event'] = team_hitters_df['PitchCall'].isin([
    'StrikeSwinging', 'StrikeCalled', 'FoulBall', 'InPlay',
    'BallCalled', 'BallIntentional', 'HitByPitch'
]) | team_hitters_df['PlayResult'].notna()
team_hitters_df = team_hitters_df[team_hitters_df['PA_event']]

# Group and count outcomes
grouped = team_hitters_df.groupby('Batter')
results = []

for batter, group in grouped:
    pos = input(f"Enter primary position for {batter}: ").strip()

    # Filter only final pitches that end a plate appearance
    plate_appearances = ((group['PitchCall'] == "InPlay") | (group['KorBB'].isin(['Strikeout', 'Walk'])) | (group['PitchCall'] == 'HitByPitch')).sum()
    
    counts = {
        '1B': (group['HitType'] == '1B').sum(),
        '2B': (group['HitType'] == '2B').sum(),
        '3B': (group['HitType'] == '3B').sum(),
        'HR': (group['HitType'] == 'HR').sum(),
        'BB': group['BB'].sum(),
        'HBP': group['HBP'].sum(),
        'PA': plate_appearances
    }
    
    # Add calculated components
    hits = counts['1B'] + counts['2B'] + counts['3B'] + counts['HR']
    total_bases = counts['1B'] + 2 * counts['2B'] + 3 * counts['3B'] + 4 * counts['HR']
    ab = counts['PA'] - counts['BB'] - counts['HBP']  # Subtract SF too if you have it

    # Calculate SLG, OBP, OPS
    slg = total_bases / ab if ab > 0 else 0
    obp = (hits + counts['BB'] + counts['HBP']) / (ab + counts['BB'] + counts['HBP']) if (ab + counts['BB'] + counts['HBP']) > 0 else 0
    ops = slg + obp



    pos_adj_table = {
        'C': 7.5, 'SS': 5, '2B': 2.5, '3B': 2.5, 'CF': 3.5,
        'LF': -2.5, 'RF': -2.5, '1B': -5, 'DH': -6
    }
    pos_adj = pos_adj_table.get(pos.upper(), 0) * (counts['PA'] / 300)
    num = sum(counts[k] * W_OBA_WEIGHTS[k] for k in W_OBA_WEIGHTS)
    denom = counts['PA']
    woba = num / denom if denom > 0 else 0
    wraa = (woba - LEAGUE_WOBA) / W_OBA_SCALE * counts['PA']
    repl = 20 * (counts['PA'] / 250)
    war = (wraa + repl + pos_adj) / RUNS_PER_WIN

    ops_plus = 100 * ((obp / LEAGUE_OBP) + ((slg / LEAGUE_SLG) - 1)) if LEAGUE_OBP > 0 and LEAGUE_SLG > 0 else 100


    total_swings = ((group['PitchCall'] == 'InPlay') | 
                (group['PitchCall'] == 'FoulBall') | 
                (group['PitchCall'] == 'FoulBallFieldable') | 
                (group['PitchCall'] == 'FoulBallNotFieldable') | 
                (group['PitchCall'] == 'StrikeSwinging')).sum()

    # Example: mark all pitches in Zones 1–9 as strikes
    good_swings = ((group['PlateLocHeight'] >= 1.524166) & 
                (group['PlateLocHeight'] <= 3.3775) &
                (group['PlateLocSide'] >= -0.830833) & 
                (group['PlateLocSide'] <= 0.830833) &
                ((group['PitchCall'] == 'InPlay') | 
                (group['PitchCall'] == 'FoulBall') | 
                (group['PitchCall'] == 'FoulBallFieldable') | 
                (group['PitchCall'] == 'FoulBallNotFieldable') | 
                (group['PitchCall'] == 'StrikeSwinging'))).sum()
    
    seager_approx = good_swings / total_swings if total_swings > 0 else 0



    results.append({
        'Batter': batter,
        'Position': pos,
        'PA': counts['PA'],
        'SLG': round(slg, 3),
        'OBP': round(obp, 3),
        'OPS': round(ops, 3),
        'wOBA': round(woba, 3),
        'wRAA': round(wraa, 2),
        'Rep. Runs': round(repl, 2),
        'Pos. Adj.': round(pos_adj, 2),
        'SEAGER': round(seager_approx, 3),
        'OPS+': int(round(ops_plus)),
        'WAR': round(war, 2)
    })

# Convert to DataFrame and display
war_df = pd.DataFrame(results)
war_df = war_df.sort_values(by='WAR', ascending=False)
#print(war_df)

def export_to_pptx(df_results, template_path, output_path):
    df_results = df_results.sort_values(by='WAR', ascending=False).reset_index(drop=True)
    # Format batter names to "Last F."
    if 'Batter' in df_results.columns:
        def _format_last_first_initial(name: str) -> str:
            try:
                parts = str(name).split(',')
                last = parts[0].strip()
                first_initial = parts[1].strip()[0] + '.' if len(parts) > 1 and len(parts[1].strip()) > 0 else ''
                if first_initial:
                    return f"{last}, {first_initial}".strip()
                return last
            except Exception:
                return str(name)
        df_results['Batter'] = df_results['Batter'].astype(str).apply(_format_last_first_initial)

    # Create a display copy with fixed decimal formatting to preserve trailing zeros
    df_display = df_results.copy()
    three_dec_cols = ['SLG', 'OBP', 'OPS', 'wOBA', 'SEAGER']
    two_dec_cols = ['wRAA', 'Rep. Runs', 'Replacement Runs', 'Pos. Adj.', 'Position Adjusment', 'WAR']
    for col in three_dec_cols:
        if col in df_display.columns:
            df_display[col] = df_display[col].apply(lambda x: f"{x:.3f}" if pd.notna(x) else "")
    for col in two_dec_cols:
        if col in df_display.columns:
            df_display[col] = df_display[col].apply(lambda x: f"{x:.2f}" if pd.notna(x) else "")
    # Load template
    prs = Presentation(template_path)
    # Use the first slide in the template
    slide = prs.slides[0]

    # Compute table size based on actual slide size; leave small margins
    emu_per_inch = 914400
    slide_width_inches = prs.slide_width / emu_per_inch
    slide_height_inches = prs.slide_height / emu_per_inch

    margin_h = 0.25
    margin_v = 0.25
    table_width = slide_width_inches - (2 * margin_h)
    table_height = slide_height_inches - (2 * margin_v)
    left = Inches(margin_h)
    top = Inches(margin_v)

    # Table dimensions
    rows, cols = df_display.shape[0] + 1, df_display.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, Inches(table_width), Inches(table_height)).table

    # Auto-adjust column widths evenly
    col_width = table_width / cols
    for i in range(cols):
        table.columns[i].width = Inches(col_width)

    # Header row
    for j, col_name in enumerate(df_display.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.name = 'Bahnschrift'
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.size = Pt(9)
        para.font.bold = True

    # Data rows
    row_height = table_height / rows  # evenly distribute row height
    for row_idx, (_, row) in enumerate(df_display.iterrows()):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.name = 'Bahnschrift'
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.font.size = Pt(9)
        table.rows[row_idx + 1].height = Inches(row_height)

    # Try to save; handle Windows file-lock by removing existing or choosing unique name
    actual_output_path = output_path
    try:
        prs.save(actual_output_path)
    except PermissionError:
        # Try to remove existing file if not in use
        try:
            if os.path.exists(actual_output_path):
                os.remove(actual_output_path)
            prs.save(actual_output_path)
        except Exception:
            # Fallback to unique filename
            base, ext = os.path.splitext(output_path)
            counter = 1
            while True:
                candidate = f"{base}-{counter}{ext}"
                if not os.path.exists(candidate):
                    prs.save(candidate)
                    actual_output_path = candidate
                    break
                counter += 1

    return actual_output_path

def pptx_to_pdf(inputFileName, outputFileName):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, 32)  # 32 = PDF
    deck.Close()
    powerpoint.Quit()
    print(f"Converted {inputFileName} → {outputFileName}")


# Export to PPTX and capture actual filename
pptx_path = export_to_pptx(war_df, "offense-template.pptx", "offense-results.pptx")

# Convert PPTX → PDF (align name if unique filename was used)
pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
pptx_to_pdf(pptx_path, pdf_path)

#PAs in 99, Runs per win, woba scale